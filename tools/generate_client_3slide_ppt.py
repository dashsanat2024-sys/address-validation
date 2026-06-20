#!/usr/bin/env python3
"""Generate 3-slide client proposal PPT in Arthavi style."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "CLIENT_3_SLIDE_PROPOSAL.pptx"

# Match existing deck palette
TEAL = RGBColor(0x02, 0x80, 0x90)
DEEP_BLUE = RGBColor(0x24, 0x36, 0x70)
LIGHT_BG = RGBColor(0xEE, 0xF4, 0xFB)
LIGHTER_BG = RGBColor(0xF4, 0xF8, 0xFB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x1A, 0x2B, 0x5E)
TEXT = RGBColor(0x1A, 0x1A, 0x1A)
MUTED = RGBColor(0x5B, 0x64, 0x73)


def box(slide, l, t, w, h, fill=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.line.fill.background()
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    return s


def rect(slide, l, t, w, h, fill=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.line.fill.background()
    if fill:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    return s


def text(slide, l, t, w, h, value, size=12, bold=False, color=TEXT, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = value
    p.alignment = align
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    return tb


def bullets(slide, l, t, w, h, items, size=10, color=TEXT):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(size)
        p.font.color.rgb = color
    return tb


def add_left_bar(slide):
    rect(slide, Inches(0), Inches(0), Inches(0.18), Inches(5.625), TEAL)


def slide_1(prs: Presentation):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_left_bar(s)
    text(s, Inches(0.55), Inches(0.25), Inches(9.2), Inches(0.55),
         "Executive Proposal", 16, True, DARK)
    text(s, Inches(0.55), Inches(0.75), Inches(9.2), Inches(0.45),
         "SAP CPI One-by-One Address Validation using Arthavi LLM", 22, True, DEEP_BLUE)

    box(s, Inches(0.55), Inches(1.35), Inches(4.45), Inches(1.55), LIGHT_BG)
    text(s, Inches(0.75), Inches(1.52), Inches(4.0), Inches(0.3), "Current Pain", 13, True, DEEP_BLUE)
    bullets(s, Inches(0.75), Inches(1.82), Inches(4.1), Inches(1.0), [
        "CPI receives inconsistent free-text addresses",
        "Manual cleansing slows SAP posting and increases errors",
        "Validation APIs alone do not map to your SAP schema",
    ], 10, MUTED)

    box(s, Inches(5.15), Inches(1.35), Inches(4.35), Inches(1.55), LIGHT_BG)
    text(s, Inches(5.35), Inches(1.52), Inches(4.0), Inches(0.3), "Arthavi Solution", 13, True, DEEP_BLUE)
    bullets(s, Inches(5.35), Inches(1.82), Inches(4.0), Inches(1.0), [
        "Validate + normalize each address in real time",
        "Map into 17 SAP fields with arthavi-address",
        "Return clean JSON response to CPI per record",
    ], 10, MUTED)

    rect(s, Inches(0.55), Inches(3.08), Inches(8.95), Inches(0.04), TEAL)
    text(s, Inches(0.55), Inches(3.22), Inches(9.0), Inches(0.32), "Business Benefits", 13, True, DEEP_BLUE)
    bullets(s, Inches(0.55), Inches(3.52), Inches(9.0), Inches(1.15), [
        "Higher SAP master-data quality and fewer downstream failures",
        "30–60% reduction in manual correction effort in phase 1",
        "Continuous learning loop from saved user corrections",
        "On-prem/private deployment option for data control",
    ], 11, TEXT)

    box(s, Inches(0.55), Inches(4.85), Inches(8.95), Inches(0.55), DEEP_BLUE)
    text(s, Inches(0.72), Inches(5.02), Inches(8.6), Inches(0.25),
         "Target KPI: ≥95% field-level mapping accuracy after stabilization", 11, True, WHITE, PP_ALIGN.CENTER)


def slide_2(prs: Presentation):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_left_bar(s)
    text(s, Inches(0.55), Inches(0.25), Inches(9.2), Inches(0.55),
         "Technical Architecture", 16, True, DARK)
    text(s, Inches(0.55), Inches(0.75), Inches(9.2), Inches(0.45),
         "How the Arthavi LLM Flow Works with SAP CPI", 22, True, DEEP_BLUE)

    # Flow row
    labels = [
        ("SAP CPI", "Address record"),
        ("Arthavi API", "/api/normalize"),
        ("Validation", "Postcodes.io / Ideal"),
        ("LLM Mapping", "arthavi-address"),
        ("SAP Output", "Mapped fields"),
    ]
    x = 0.55
    for i, (h, sub) in enumerate(labels):
        fill = RGBColor(0x02, 0x79, 0x8A + min(i, 3) * 4) if i < 4 else DEEP_BLUE
        box(s, Inches(x), Inches(1.55), Inches(1.72), Inches(1.12), fill)
        text(s, Inches(x + 0.08), Inches(1.72), Inches(1.56), Inches(0.27), h, 10, True, WHITE, PP_ALIGN.CENTER)
        text(s, Inches(x + 0.08), Inches(2.0), Inches(1.56), Inches(0.45), sub, 9, False, WHITE, PP_ALIGN.CENTER)
        if i < len(labels) - 1:
            rect(s, Inches(x + 1.72), Inches(2.03), Inches(0.20), Inches(0.09), TEAL)
        x += 1.9

    box(s, Inches(0.55), Inches(3.02), Inches(4.45), Inches(1.95), LIGHTER_BG)
    text(s, Inches(0.75), Inches(3.2), Inches(4.0), Inches(0.3), "Tech Stack", 13, True, DEEP_BLUE)
    bullets(s, Inches(0.75), Inches(3.5), Inches(4.1), Inches(1.35), [
        "Python + Flask API service",
        "Ollama runtime with arthavi-address model",
        "Postcodes.io (free) or Ideal Postcodes validator",
        "SAP CPI REST integration (one-by-one)",
        "Optional correction capture + retraining pipeline",
    ], 10, MUTED)

    box(s, Inches(5.15), Inches(3.02), Inches(4.35), Inches(1.95), LIGHTER_BG)
    text(s, Inches(5.35), Inches(3.2), Inches(4.0), Inches(0.3), "Reliability & Governance", 13, True, DEEP_BLUE)
    bullets(s, Inches(5.35), Inches(3.5), Inches(4.0), Inches(1.35), [
        "Deterministic schema output for SAP posting",
        "Timeout-safe fallback handling for cold model starts",
        "Audit-friendly validation metadata and warnings",
        "Human-in-loop corrections for controlled improvement",
    ], 10, MUTED)


def slide_3(prs: Presentation):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_left_bar(s)
    text(s, Inches(0.55), Inches(0.25), Inches(9.2), Inches(0.55),
         "Commercial Proposal", 16, True, DARK)
    text(s, Inches(0.55), Inches(0.75), Inches(9.2), Inches(0.45),
         "Cost, Deployment Options, and Delivery Plan", 22, True, DEEP_BLUE)

    box(s, Inches(0.55), Inches(1.35), Inches(4.45), Inches(1.65), RGBColor(0xEC, 0xFD, 0xF5))
    text(s, Inches(0.75), Inches(1.52), Inches(4.0), Inches(0.3), "Option A (Recommended Start)", 12, True, RGBColor(0x04, 0x78, 0x57))
    bullets(s, Inches(0.75), Inches(1.82), Inches(4.1), Inches(1.08), [
        "Postcodes.io + local arthavi-address model",
        "Validation API cost: £0",
        "External LLM API cost: £0 (local inference)",
        "Best ROI for pilot and phased rollout",
    ], 10, MUTED)

    box(s, Inches(5.15), Inches(1.35), Inches(4.35), Inches(1.65), RGBColor(0xE0, 0xF4, 0xFF))
    text(s, Inches(5.35), Inches(1.52), Inches(4.0), Inches(0.3), "Option B (Higher Precision)", 12, True, RGBColor(0x03, 0x69, 0xA1))
    bullets(s, Inches(5.35), Inches(1.82), Inches(4.0), Inches(1.08), [
        "Ideal Postcodes + arthavi-address",
        "Pay-per-lookup validation (pence/record)",
        "Stronger street-level confidence for critical data",
    ], 10, MUTED)

    text(s, Inches(0.55), Inches(3.22), Inches(9.0), Inches(0.32), "Implementation Timeline", 13, True, DEEP_BLUE)
    weeks = [
        ("Week 1", "CPI contract, endpoint setup, sample payload testing"),
        ("Week 2", "UAT with live addresses, correction workflow enabled"),
        ("Week 3", "Accuracy review, retrain pass, go-live recommendation"),
    ]
    for i, (wk, desc) in enumerate(weeks):
        left = 0.55 + i * 3.0
        box(s, Inches(left), Inches(3.52), Inches(2.82), Inches(1.25), LIGHT_BG)
        text(s, Inches(left + 0.12), Inches(3.68), Inches(2.55), Inches(0.24), wk, 11, True, TEAL)
        text(s, Inches(left + 0.12), Inches(3.95), Inches(2.55), Inches(0.72), desc, 9, False, MUTED)

    box(s, Inches(0.55), Inches(4.95), Inches(8.95), Inches(0.58), DEEP_BLUE)
    text(s, Inches(0.72), Inches(5.12), Inches(8.6), Inches(0.28),
         "Proposal ask: approve 3-week pilot to validate accuracy, effort reduction, and production readiness.",
         10, True, WHITE, PP_ALIGN.CENTER)


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625)
    slide_1(prs)
    slide_2(prs)
    slide_3(prs)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    main()

