#!/usr/bin/env python3
"""Generate 2-slide client proposal PowerPoint."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "CLIENT_2_SLIDE_PROPOSAL.pptx"

BRAND = RGBColor(0x6B, 0x52, 0xB0)
BRAND_DEEP = RGBColor(0x1B, 0x1A, 0x3E)
ACCENT = RGBColor(0x04, 0xA0, 0xE8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x6B, 0x6B, 0x6B)
DARK = RGBColor(0x1A, 0x1A, 0x1A)


def _box(slide, left, top, width, height, fill=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    return shape


def _text(slide, left, top, width, height, text, size=14, bold=False, color=DARK, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return tf


def _bullets(slide, left, top, width, height, items, size=11, color=DARK, bold_prefix=True):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if ": " in item and bold_prefix:
            head, tail = item.split(": ", 1)
            p.text = f"• {head}: {tail}"
            p.font.size = Pt(size)
            p.font.color.rgb = color
            p.level = 0
        else:
            p.text = f"• {item}"
            p.font.size = Pt(size)
            p.font.color.rgb = color
    return tf


def slide1(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _box(slide, Inches(0), Inches(0), Inches(10), Inches(0.9), BRAND_DEEP)
    _text(slide, Inches(0.5), Inches(0.2), Inches(9), Inches(0.5),
          "Arthavi Address Validation LLM", 12, False, WHITE)
    _text(slide, Inches(0.5), Inches(0.45), Inches(9), Inches(0.45),
          "SAP CPI Address Validation — Record-by-Record", 22, True, WHITE)

    _text(slide, Inches(0.5), Inches(1.05), Inches(4.5), Inches(0.35),
          "The Challenge", 14, True, BRAND_DEEP)
    _bullets(slide, Inches(0.5), Inches(1.35), Inches(4.5), Inches(1.5), [
        "SAP CPI sends addresses one-by-one from multiple sources",
        "Messy formats: reordered fields, abbreviations, mixed casing",
        "Manual cleanup slows SAP posting and causes errors",
        "Postcode APIs validate but don't map to your SAP schema",
    ], size=10)

    _text(slide, Inches(5.2), Inches(1.05), Inches(4.3), Inches(0.35),
          "Our Solution", 14, True, BRAND_DEEP)
    _bullets(slide, Inches(5.2), Inches(1.35), Inches(4.3), Inches(1.5), [
        "Existing Arthavi app exposed as API behind SAP CPI",
        "Validate postcode → normalize with arthavi-address LLM",
        "Return 17-field SAP client structure per record",
        "Corrections feed retraining for continuous improvement",
    ], size=10)

    _box(slide, Inches(0.5), Inches(2.95), Inches(9), Inches(0.04), ACCENT)

    _text(slide, Inches(0.5), Inches(3.1), Inches(4.5), Inches(0.3),
          "How the LLM Helps", 13, True, BRAND)
    _bullets(slide, Inches(0.5), Inches(3.4), Inches(4.5), Inches(1.2), [
        "Handles reordered UK addresses better than rules alone",
        "Maps co, street_2–5, district, city, postal_code correctly",
        "Learns vendor-specific patterns from your corrections",
        "Consistent output for CPI sequential processing",
    ], size=10, bold_prefix=False)

    _text(slide, Inches(5.2), Inches(3.1), Inches(4.3), Inches(0.3),
          "Client Benefits & KPIs", 13, True, BRAND)
    _bullets(slide, Inches(5.2), Inches(3.4), Inches(4.3), Inches(1.2), [
        "Higher data quality: fewer wrong SAP postings",
        "Faster turnaround: less manual address cleanup",
        "Data control: on-prem / private deployment",
        "Target: ≥95% field accuracy after stabilization",
        "30–60% reduction in manual corrections (phase 1)",
    ], size=10, bold_prefix=False)

    _box(slide, Inches(0.5), Inches(5.0), Inches(9), Inches(0.55), RGBColor(0xF0, 0xEA, 0xFF))
    _text(slide, Inches(0.65), Inches(5.1), Inches(8.7), Inches(0.4),
          "Not a generic chatbot — a domain-tuned mapping engine for SAP address structures.",
          11, True, BRAND_DEEP, PP_ALIGN.CENTER)


def slide2(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _box(slide, Inches(0), Inches(0), Inches(10), Inches(0.9), BRAND_DEEP)
    _text(slide, Inches(0.5), Inches(0.2), Inches(9), Inches(0.5),
          "Arthavi", 12, False, WHITE)
    _text(slide, Inches(0.5), Inches(0.45), Inches(9), Inches(0.45),
          "Architecture, Cost & Implementation", 22, True, WHITE)

    _text(slide, Inches(0.5), Inches(1.05), Inches(4.8), Inches(0.3),
          "End-to-End Flow (SAP CPI)", 13, True, BRAND_DEEP)
    flow = (
        "SAP CPI  →  /api/normalize  →  Preprocess\n"
        "    ↓\n"
        "Postcodes.io (validate)  →  arthavi-address LLM\n"
        "    ↓\n"
        "SAP mapped fields + metadata  →  CPI response\n"
        "    ↓\n"
        "Corrections  →  Retrain  →  Improved model"
    )
    _text(slide, Inches(0.5), Inches(1.35), Inches(4.8), Inches(1.8), flow, 10, False, DARK)

    _text(slide, Inches(5.4), Inches(1.05), Inches(4.1), Inches(0.3),
          "Tech Stack", 13, True, BRAND_DEEP)
    _bullets(slide, Inches(5.4), Inches(1.35), Inches(4.1), Inches(1.5), [
        "API: Python + Flask",
        "LLM: Ollama + arthavi-address (fine-tuned)",
        "Validation: Postcodes.io (free) or Ideal Postcodes",
        "Integration: SAP CPI REST per address",
        "Deploy: on-prem VM or Cloud Run + Ollama VM",
    ], size=10, bold_prefix=False)

    _box(slide, Inches(0.5), Inches(3.2), Inches(4.3), Inches(1.35), RGBColor(0xEC, 0xFD, 0xF5))
    _text(slide, Inches(0.65), Inches(3.28), Inches(4), Inches(0.25),
          "Option A — Recommended Start", 11, True, RGBColor(0x04, 0x78, 0x57))
    _bullets(slide, Inches(0.65), Inches(3.55), Inches(4), Inches(0.95), [
        "Postcodes.io + local arthavi-address",
        "Validation cost: £0",
        "LLM cost: £0 external (local inference)",
        "Best for pilot & fast ROI",
    ], size=9, bold_prefix=False)

    _box(slide, Inches(5.0), Inches(3.2), Inches(4.5), Inches(1.35), RGBColor(0xE0, 0xF4, 0xFF))
    _text(slide, Inches(5.15), Inches(3.28), Inches(4.2), Inches(0.25),
          "Option B — Premium Validation", 11, True, RGBColor(0x03, 0x69, 0xA1))
    _bullets(slide, Inches(5.15), Inches(3.55), Inches(4.2), Inches(0.95), [
        "Ideal Postcodes + arthavi-address",
        "Pay-per-lookup (pence per record)",
        "Street-level confidence for delivery-critical data",
    ], size=9, bold_prefix=False)

    _text(slide, Inches(0.5), Inches(4.75), Inches(9), Inches(0.3),
          "3-Week Delivery Plan", 13, True, BRAND_DEEP)
    cols = [
        ("Week 1", "CPI endpoint contract\n+ sample payload testing"),
        ("Week 2", "UAT with real addresses\n+ correction loop"),
        ("Week 3", "Accuracy review\n+ retrain + go-live"),
    ]
    for i, (title, body) in enumerate(cols):
        left = Inches(0.5 + i * 3.15)
        _box(slide, left, Inches(5.05), Inches(2.95), Inches(0.9), RGBColor(0xF5, 0xF7, 0xFB))
        _text(slide, left + Inches(0.1), Inches(5.1), Inches(2.75), Inches(0.25), title, 11, True, BRAND)
        _text(slide, left + Inches(0.1), Inches(5.35), Inches(2.75), Inches(0.55), body, 9, False, MUTED)

    _text(slide, Inches(0.5), Inches(6.05), Inches(9), Inches(0.3),
          "Low-risk adoption: app already built — this extends it for SAP CPI, not a re-platform.",
          10, False, MUTED, PP_ALIGN.CENTER)


def main() -> None:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    slide1(prs)
    slide2(prs)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    main()
