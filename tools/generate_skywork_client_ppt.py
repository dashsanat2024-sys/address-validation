#!/usr/bin/env python3
"""Generate Skywork-style 2-slide SAP CPI client proposal PowerPoint.

Uses the bundled Skywork template (card layout, flow diagram, KPI tiles).
The reference design matches skywork.ai output: widescreen 16:9, navy/teal
cards, before/after address mapping visual, architecture + cost + 3-week plan.

Usage:
    python tools/generate_skywork_client_ppt.py
    python tools/generate_skywork_client_ppt.py --out docs/my_deck.pptx
"""

from __future__ import annotations

import argparse
import shutil
from datetime import datetime
from pathlib import Path

from pptx import Presentation

ROOT = Path(__file__).resolve().parents[1]
TEMPLATE = ROOT / "docs" / "ppt_templates" / "skywork_sap_cpi_2slide.pptx"
DEFAULT_OUT = ROOT / "docs" / "SAP_CPI_Address_Validation_Arthavi_LLM.pptx"

# Text shape names -> replacement values (optional overrides)
TEXT_UPDATES: dict[str, str] = {}


def _update_proposal_date(prs: Presentation, month_year: str) -> None:
    """Refresh CLIENT PROPOSAL date stamp on both slides."""
    label = month_year.upper()
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            if not text.startswith("CLIENT PROPOSAL"):
                continue
            sep = " · " if "·" in text else " - "
            shape.text_frame.text = f"CLIENT PROPOSAL{sep}{label}"


def _apply_text_overrides(prs: Presentation, overrides: dict[str, str]) -> None:
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.name in overrides and shape.has_text_frame:
                shape.text_frame.text = overrides[shape.name]


def generate(
    out_path: Path,
    *,
    template: Path = TEMPLATE,
    month_year: str | None = None,
    text_overrides: dict[str, str] | None = None,
) -> Path:
    if not template.is_file():
        raise FileNotFoundError(
            f"Skywork template not found: {template}\n"
            "Copy the reference deck to docs/ppt_templates/skywork_sap_cpi_2slide.pptx"
        )

    out_path.parent.mkdir(parents=True, exist_ok=True)
    shutil.copy2(template, out_path)

    prs = Presentation(str(out_path))
    if month_year:
        _update_proposal_date(prs, month_year)
    if text_overrides:
        _apply_text_overrides(prs, text_overrides)
    prs.save(str(out_path))
    return out_path


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument(
        "--out",
        type=Path,
        default=DEFAULT_OUT,
        help=f"Output path (default: {DEFAULT_OUT.relative_to(ROOT)})",
    )
    parser.add_argument(
        "--template",
        type=Path,
        default=TEMPLATE,
        help="Skywork template .pptx",
    )
    parser.add_argument(
        "--month-year",
        default=datetime.now().strftime("%B %Y"),
        help="Proposal date label, e.g. 'June 2026'",
    )
    parser.add_argument(
        "--no-date-update",
        action="store_true",
        help="Keep template date stamp unchanged",
    )
    args = parser.parse_args()

    path = generate(
        args.out,
        template=args.template,
        month_year=None if args.no_date_update else args.month_year,
        text_overrides=TEXT_UPDATES or None,
    )
    print(f"Wrote {path} ({path.stat().st_size // 1024} KB)")


if __name__ == "__main__":
    main()
